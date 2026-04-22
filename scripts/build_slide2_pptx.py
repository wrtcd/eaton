"""
PPT: (1) TEMPO FOR, (2) reprojected grid, (3) column → mass, (4) plume map + max/min+,
(5) plume mask f_p min / max and order-of-magnitude difference.
Run: py -3 scripts/build_slide2_pptx.py
Regenerate the figure: py -3 scripts/plot_max_mass_zoom_inset.py
"""
from __future__ import annotations

import sys
import math
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

_ROOT = Path(__file__).resolve().parents[1]
_OUT = _ROOT / "step-by-step" / "presentations" / "slide2_column_to_mass.pptx"
_DEFAULT_GRID = _ROOT / "step-by-step" / "08 mass" / "tempo_mass_no2_kg_per_pixel.tif"
_DEFAULT_FP = _ROOT / "step-by-step" / "04 plume fraction" / "tempo_fp_plume_utm11_clipped.tif"
_PLUME_FIG = _ROOT / "step-by-step" / "presentations" / "max_mass_pixel_inset.png"
_OUT.parent.mkdir(parents=True, exist_ok=True)

S = {
    "0": "\u2070", "1": "\u00b9", "2": "\u00b2", "3": "\u00b3", "4": "\u2074", "5": "\u2075",
    "6": "\u2076", "7": "\u2077", "8": "\u2078", "9": "\u2079", "-": "\u207b",
}
SUPMIN = S["-"] + S["1"]  # ⁻¹


def T(exp: str) -> str:
    return "10" + "".join(S[c] for c in exp)


def grid_from_geotiff(path: Path) -> tuple[float, float, float, str] | None:
    """E–W m, N–S m, area m², path str. None if file missing or error."""
    try:
        import rasterio
    except ImportError:
        return None
    if not path.is_file():
        return None
    with rasterio.open(path) as ds:
        t = ds.transform
        a, b, d, e = t.a, t.b, t.d, t.e
        ew = float((a**2 + b**2) ** 0.5)  # column direction (E–W, easting)
        ns = float((d**2 + e**2) ** 0.5)  # row direction (N–S, northing, length)
        area_m2 = float(abs(a * e - b * d))
    return ew, ns, area_m2, str(path.name)


def plume_mask_fmin_fmax(
    path: Path,
) -> tuple[float, float, float, str] | None:
    """
    Min / max f_p over positive, in-range pixels; log10(max/min); filename.
    None if file missing, import error, or no valid plume cells.
    """
    try:
        import numpy as np
        import rasterio
    except ImportError:
        return None
    if not path.is_file():
        return None
    with rasterio.open(path) as ds:
        a = np.asarray(ds.read(1), dtype=np.float64)
        nod = ds.nodata
    m = np.isfinite(a) & (a > 0.0) & (a <= 1.0)
    if nod is not None and np.isfinite(nod):
        m &= a != float(nod)
    if not np.any(m):
        return None
    f_lo = float(np.min(a[m]))
    f_hi = float(np.max(a[m]))
    if f_lo <= 0.0 or not math.isfinite(f_hi / f_lo):
        return None
    log_orders = math.log10(f_hi / f_lo)
    return f_lo, f_hi, log_orders, str(path.name)


def main() -> int:
    f_p, n_a = "f\u209a", "N\u2090"
    cm2_area = "cm" + S["2"]
    per_cm2 = "cm" + S["-"] + S["2"]
    m2 = "m" + S["2"]
    km2 = "km" + S["2"]
    m_n_o2 = "M" + "N" + "O" + "\u2082"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def p_par(tf, text, siz, bold=False, sp=3):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(siz)
        p.font.bold = bold
        p.space_after = Pt(sp)

    # —— Slide 1: TEMPO center
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s1.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(6.9))
    tf = tb.text_frame
    tf.clear()
    tf.paragraphs[0].text = "TEMPO: ground pixel, FOR center"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].space_after = Pt(8)

    p_par(tf, "N–S: 2.00 km          E–W: 4.75 km", 22, True, 4)
    p_par(tf, f"Area = 2.00 km × 4.75 km  =  9.50 {km2}  =  9.50×{T('6')}  {m2}", 20, True, 4)
    p_par(tf, f"2.00 km = 2.00×{T('3')} m          4.75 km = 4.75×{T('3')} m", 18, False, 2)
    p_par(
        tf,
        f"2.00×{T('3')} m × 4.75×{T('3')} m  =  9.50×{T('6')}  {m2}",
        18,
        False,
        2,
    )
    p_par(
        tf,
        f"1  {m2}  =  {T('4')}  {cm2_area}  →  9.50×{T('10')}  {cm2_area}",
        20,
        True,
        4,
    )
    p_par(tf, "tempo.si.edu/instrument  ·  33.5°N, 89.2°W", 12, False, 0)

    # —— Slide 2: reprojected study grid (from affine)
    s_grid = prs.slides.add_slide(prs.slide_layouts[6])
    tb3 = s_grid.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(6.9))
    tg = tb3.text_frame
    tg.clear()
    g = grid_from_geotiff(_DEFAULT_GRID)
    tg.paragraphs[0].text = "Reprojected study grid (this project)"
    tg.paragraphs[0].font.size = Pt(32)
    tg.paragraphs[0].font.bold = True
    tg.paragraphs[0].space_after = Pt(6)

    if g is not None:
        ew_m, ns_m, area_m2, fname = g
        ew_km, ns_km = ew_m / 1000, ns_m / 1000
        area_km2 = area_m2 / 1.0e6
        area_cm2 = area_m2 * 1.0e4
        p_par(
            tg,
            f"GeoTIFF: …/{fname}  ·  EPSG:32611  ·  (easting × northing from affine)",
            12,
            False,
            6,
        )
        p_par(
            tg,
            f"E–W (pixel, x):  {ew_m:,.2f} m  =  {ew_km:.4f} km",
            20,
            True,
            2,
        )
        p_par(
            tg,
            f"N–S (pixel, y):  {ns_m:,.2f} m  =  {ns_km:.4f} km",
            20,
            True,
            4,
        )
        p_par(
            tg,
            f"Area  =  E–W × N–S  =  {ew_m:.2f} m × {ns_m:.2f} m  =  {area_m2:,.2f}  {m2}",
            19,
            True,
            2,
        )
        p_par(
            tg,
            f"=  {area_km2:.4f}  {km2}  =  {area_m2/1.0e7:.2f}×{T('7')}  {m2}",
            18,
            False,
            2,
        )
        p_par(
            tg,
            f"1  {m2}  =  {T('4')}  {cm2_area}  →  {area_cm2/1.0e11:.4f}×{T('11')}  {cm2_area}",
            20,
            True,
            4,
        )
    else:
        p_par(
            tg,
            f"Run failed: add rasterio and file {_DEFAULT_GRID.name}",
            14,
            False,
            0,
        )

    # —— Slide 3: column → mass
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb2 = s2.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(6.9))
    t2 = tb2.text_frame
    t2.clear()
    t2.paragraphs[0].text = "Column → mass"
    t2.paragraphs[0].font.size = Pt(32)
    t2.paragraphs[0].font.bold = True
    t2.paragraphs[0].space_after = Pt(4)

    p_par(t2, "Symbols  (TEMPO grid + plume product)", 15, True, 2)
    p_par(
        t2,
        f"σ   —  plume column  ({f_p}·ΔVCD);  same kind of field as TEMPO VCD  (molec·{per_cm2})",
        13,
        False,
        0,
    )
    p_par(
        t2,
        f"{f_p}  —  mean Planet plume mask in that TEMPO cell  (0…1)  (step 04)",
        13,
        False,
        0,
    )
    p_par(
        t2,
        f"ΔVCD  —  VCD  −  VCD_bg  (excess column after local background; step 07)",
        13,
        False,
        0,
    )
    p_par(
        t2,
        f"A  —  ground footprint of one raster cell  ({m2}  →  {cm2_area}; reprojected GeoTIFF)",
        13,
        False,
        0,
    )
    p_par(
        t2,
        f"{n_a}  —  Avogadro.        {m_n_o2}  —  molar mass of  NO\u2082  (kg·mol{SUPMIN})",
        13,
        False,
        6,
    )

    p_par(t2, f"σ = {f_p}·ΔVCD  (molec·{per_cm2})", 22, False, 4)
    p_par(t2, f"m  (kg)  =  σ  ·  (A  in  {cm2_area})  /  ({n_a}  ·  {m_n_o2})", 24, True, 6)
    p_par(
        t2,
        f"{n_a} = 6.02214076×{T('23')} mol{SUPMIN}     {m_n_o2} = 0.0460055 kg·mol{SUPMIN}",
        15,
        False,
        4,
    )
    p_par(t2, f"e.g. max-mass cell:  σ  ≈  9.1×{T('16')}  molec·{per_cm2}", 14, False, 0)

    # —— Slide 4: reprojected plume — max mass vs smallest positive m (figure)
    s_map = prs.slides.add_slide(prs.slide_layouts[6])
    tb4 = s_map.shapes.add_textbox(Inches(0.65), Inches(0.38), Inches(12.1), Inches(1.05))
    t4 = tb4.text_frame
    t4.clear()
    t4.paragraphs[0].text = "Reprojected grid: where are the max and weakest positive cells?"
    t4.paragraphs[0].font.size = Pt(28)
    t4.paragraphs[0].font.bold = True
    t4.paragraphs[0].space_after = Pt(2)
    p4 = t4.add_paragraph()
    p4.text = (
        "Overview (WGS 84, kg NO\u2082 per cell) + zooms: same pixel area A; "
        "colors scale with \u03c3 = f_p·\u0394VCD. Arrows: overview markers \u2192 insets."
    )
    p4.font.size = Pt(12)
    p4.space_after = Pt(0)

    fig_path = _PLUME_FIG
    if fig_path.is_file():
        s_map.shapes.add_picture(
            str(fig_path.resolve()),
            Inches(0.45),
            Inches(1.22),
            width=Inches(12.4),
        )
    else:
        miss = s_map.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11.5), Inches(2.0))
        mt = miss.text_frame
        mt.text = f"Add figure by running:  python scripts/plot_max_mass_zoom_inset.py  →  {fig_path.name}"
        mt.paragraphs[0].font.size = Pt(14)

    # —— Slide 5: plume mask f_p — min, max, order-of-magnitude difference
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    tb5 = s5.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(6.5))
    t5 = tb5.text_frame
    t5.clear()
    t5.paragraphs[0].text = "Plume mask  (Planet fraction on TEMPO grid)"
    t5.paragraphs[0].font.size = Pt(30)
    t5.paragraphs[0].font.bold = True
    t5.paragraphs[0].space_after = Pt(8)

    fp_stats = plume_mask_fmin_fmax(_DEFAULT_FP)
    if fp_stats is not None:
        f_lo, f_hi, log_o, fn = fp_stats
        ratio = f_hi / f_lo
        p_par(
            t5,
            f"{f_p}  from  …/{fn}  (step 04).  Range over pixels with  0 < {f_p} \u2264 1  (screened nodata removed).",
            14,
            False,
            10,
        )
        p_par(
            t5,
            f"min  {f_p}  =  {f_lo:.4g}                    max  {f_p}  =  {f_hi:.4g}",
            24,
            True,
            12,
        )
        p_par(
            t5,
            f"max / min  =  {ratio:.4g}  \u2192  about  {log_o:.1f}  orders of magnitude  (log\u2081\u2080 of the ratio).",
            20,
            False,
            0,
        )
    else:
        p_par(
            t5,
            f"Add plume-fraction grid:  …/step-by-step/04 plume fraction/{_DEFAULT_FP.name}",
            16,
            False,
            6,
        )
        p_par(
            t5,
            f"This slide will list  min / max  {f_p}  and  log\u2081\u2080(max/min)  when that file is on disk next to the project.",
            15,
            False,
            0,
        )

    out = str(_OUT)
    try:
        prs.save(out)
    except OSError:
        alt = out.replace(".pptx", "_new.pptx")
        prs.save(alt)
        print("Wrote (original locked)", alt, "5 slides")
    else:
        print("Wrote", out, "5 slides")
    if g:
        print("Grid:", f"E–W {g[0]:.2f} m, N–S {g[1]:.2f} m, A = {g[2]:.2f} m²")
    if fp_stats:
        flo, fhi, lo, _fn = fp_stats
        print("f_p plume mask:", f"min={flo:.4g} max={fhi:.4g} log10(max/min)={lo:.2f}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
